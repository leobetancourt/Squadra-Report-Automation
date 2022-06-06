from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from pptx.enum.chart import XL_TICK_MARK

from random import random

import json

TRAITS_ENGLISH = [
    ['Enthusiastic', 'Handles Conflict', 'Optimistic', 'Persistent'],
    ['Collaborative', 'Enlists Cooperation', 'Helpful', 'Organizational Compatibility'],
    ['Assertive', 'Warmth / empathy', 'Diplomatic', 'Outgoing', 'Frank'],
    ['Risking', 'Experimenting', 'Innovative', 'Wants Challenge'],
    ['Open / reflective', 'Analyzes Pitfalls', 'Analytical', 'Intuitive', 'Precise', 'Systematic'],
    ['Authoritative', 'Certain', 'Coaching', 'Wants To Lead', 'Effective Enforcing', 'Flexible', 'Influencing', 'Receives Correction', 'Pressure Tolerance', 'Cause Motivated'],
    ['Doesn\'t Need Structure', 'Organized', 'Planning', 'Tolerance Of Structure'],
    ['Self-acceptance', 'Self-improvement', 'Comfort With Conflict', 'Manages Stress Well', 'Relaxed', 'Self-motivated', 'Takes Initiative', 'Wants Development', 'Warmth / empathy']
]

TRAITS_SPANISH = [
    ['Entusiasta', 'Maneja conflictos', 'Optimista', 'Persistente'],
    ['Colaborador/a', 'Cooperativo/a', 'Servicial', 'Compatibilidad organizativa'],
    ['Asertivo/a', 'Cálido/a / Empático/a', 'Diplomático/a', 'Extravertido/a', 'Franco/a'],
    ['Arriesgador/a', 'Experimentador/a', 'Innovador/a', 'Desea retos o desafíos'],
    ['Abierto/a', 'Analiza los fracasos', 'Analítico/a', 'Intuitivo/a', 'Preciso/a', 'Sistemático/a'],
    ['Autoritorio/a', 'Certero', 'Coaching', 'Desea liderar', 'Eficaz para hacer cumplir', 'Flexible', 'Influyente', 'Recibe correciones', 'Tolerancia a la presión', 'Motivado por una causa'],
    ['No necesita estructura', 'Organizador/a', 'Plaificador/a', 'Tolerancia de estructura'],
    ['Autoaceptación', 'Autosuperación', 'Confort con el conflicto', 'Gestiona bien el estrés', 'Relajado/a', 'Automotivado/a', 'Lleva la iniciativa', 'Desarrollo personal', 'Cálido/a / Empático/a'],
]

def formatTitles(title):
    result = ''
    for i in range(len(title)-1):
        if title[i] != ' ':
            result += title[i]
        elif title[i] == ' ' and title[i+1] == ' ' and title[i+2] == ' ':
            result += title[i]

    result += title[-1]
    result = result.lower()
    result = result.capitalize()
    return result

def mapNum(value, leftMin, leftMax, rightMin, rightMax):
    # Figure out how 'wide' each range is
    leftSpan = leftMax - leftMin
    rightSpan = rightMax - rightMin

    # Convert the left range into a 0-1 range (float)
    valueScaled = float(value - leftMin) / float(leftSpan)

    # Convert the 0-1 range into a value in the right range.
    return rightMin + (valueScaled * rightSpan)

def openPresentation(file_name):
    prs = Presentation(file_name)
    return prs

def slide1(prs, n):
    slide1 = prs.slides[0]

    for shape in slide1.shapes:
        if round(shape.top.cm, 2) == 9.16:
            name = shape

    text_frame = name.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Usuario'
    font = run.font
    font.name = 'Raleway'
    font.bold = True
    font.size = Pt(16)
    font.color.rgb = RGBColor(117, 62, 255)
    p = text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = n
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(16)
    font.color.rgb = RGBColor(0, 0, 0)

def updateTraitSlide(prs, i, traitArr, scores, compScores):
    slide = prs.slides[i]
    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = traitArr
    series = chart_data.add_series('', scores)

    # add chart to slide --------------------
    x, y, cx, cy = Cm(1.44), Cm(2.81), Cm(9.65), Cm(6.49)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(117, 62, 255)
    chart.has_title = False
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.tick_labels.font.size = Pt(7)
    category_axis.tick_labels.font.name = 'Raleway'
    category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)
    
    value_axis = chart.value_axis
    value_axis.maximum_scale = 10.0
    value_axis.has_minor_gridlines = True
    value_axis.has_major_gridlines = False
    value_axis.minor_unit = 1
    value_axis.major_unit = 1
    value_axis.tick_labels.font.size = Pt(7)
    value_axis.tick_labels.font.name = 'Raleway'
    value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    title = ''
    for shape in slide.shapes:
        if round(shape.top.cm, 2) == 11.6:
            avgBar = shape
        if round(shape.top.cm, 2) == 10.41:
            desarrollo = shape
        if round(shape.top.cm, 2) == 11.18:
            avgScoreText = shape
        if round(shape.top.cm, 2) == 0.14:
            title = shape.text_frame.paragraphs[0].runs[0].text
            title = formatTitles(title)

    avgScore = sum(scores) / len(scores)
    w = mapNum(avgScore, 0, 10, 0, 6.07)
    avgBar.width = Cm(w)

    compScores.append({"competencia": title, "score": avgScore})

    text_frame = avgScoreText.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(round(avgScore, 1))
    font = run.font
    font.name = 'Raleway'
    font.bold = True
    font.size = Pt(16)
    font.color.rgb = RGBColor(0, 0, 0)

    text_frame = desarrollo.text_frame
    p = text_frame.add_paragraph()
    run = p.add_run()

    if avgScore >= 8.5:
        run.text = 'Bien Desarrollado'
    elif avgScore >= 7:
        run.text = 'Desarrollado'
    elif avgScore >= 5.5:
        run.text = 'Moderadamente Desarrollado'
    else:
        run.text = 'Necesita Desarrollo'

    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(16)
    font.color.rgb = RGBColor(0, 0, 0)

def slide11(prs, compScores):
    slide = prs.slides[10]

    for shape in slide.shapes:
        if round(shape.left.cm, 2) == 1.96:
            highestShape = shape
        if round(shape.left.cm, 2) == 13.47:
            lowestShape = shape

    text_frame = highestShape.text_frame
    p = text_frame.paragraphs[1]
    p.level = 1
    run = p.add_run()
    run.text = compScores[0]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    p = text_frame.paragraphs[2]
    p.level = 1
    run = p.add_run()
    run.text = compScores[1]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)
    
    p = text_frame.paragraphs[3]
    p.level = 1
    run = p.add_run()
    run.text = compScores[2]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    text_frame = lowestShape.text_frame
    p = text_frame.paragraphs[1]
    p.level = 1
    run = p.add_run()
    run.text = compScores[-1]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    p = text_frame.paragraphs[2]
    p.level = 1
    run = p.add_run()
    run.text = compScores[-2]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    p = text_frame.paragraphs[3]
    p.level = 1
    run = p.add_run()
    run.text = compScores[-3]['competencia']
    font = run.font
    font.name = 'Raleway'
    font.bold = False
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

def main():
    with open('./json_data/30052022.json') as dataFile:
        data = json.load(dataFile)
        for user in data:
            prs = openPresentation('template.pptx')
            competenciaScores = []
            if len(user['txt_data']['harrison']) > 0:
                slide1(prs, user['name'])
                for i in range(len(TRAITS_ENGLISH)):
                    # get trait scores for this user from json
                    traitScores = user['txt_data']['harrison']
                    scores = []
                    
                    for j in range(len(TRAITS_ENGLISH[i])):
                        for scoreObject in traitScores:
                            if scoreObject['trait'] == TRAITS_ENGLISH[i][j]:
                                scores.append(scoreObject['score'])
                    
                    # update corresponding slide on presentation
                    updateTraitSlide(prs, i+2, TRAITS_SPANISH[i], scores, competenciaScores)

                def func(e):
                    return e['score']
                
                competenciaScores.sort(reverse=True, key=func)

                slide11(prs, competenciaScores)
                
                prs.save('./30052022/' + user['name'] + '.pptx')
                print('Report generated for ' + user['name'])

if __name__ == '__main__':
    print('running')
    main()